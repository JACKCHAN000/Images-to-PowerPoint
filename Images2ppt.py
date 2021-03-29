# -*- coding: UTF-8 -*-
from pptx import Presentation
from pptx.util import Inches
import glob
format = ".jpg"

prs = Presentation()
# create pptx
lst_imgs = [i for i in glob.glob("*"+format)]
# get images on the folder

for img in lst_imgs:
    name = img.split(format)[0]
    title_slide_layout = prs.slide_layouts[1]
    # set slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    # crete new page
    title = slide.shapes.title
    title.text = name
    # set tittle
    left = Inches(0)
    height = Inches(5)
    left = Inches(0)
    top = Inches(3)
    pic = slide.shapes.add_picture(img, left, top, height=height)
    # add images


prs.save("sample.pptx")
# save ppt
print("done")
